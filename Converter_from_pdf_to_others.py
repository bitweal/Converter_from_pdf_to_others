from pdf2docx import Converter
import pdfplumber
import pandas as pd
import pdfplumber
from pptx import Presentation
from pptx.util import Inches
import os
import fitz
from PIL import Image


def pdf_to_docx(pdf_file, output_file, page):
    cv = Converter(pdf_file)
    try:
        if page is None:   
            start_page, end_page = 0, None
            cv.convert(output_file, start=start_page, end=end_page)
        elif len(page) == 2:
            start_page, end_page = page[0], page[1]
            cv.convert(output_file, start=start_page, end=end_page)
        else:          
            cv.convert(output_file, pages=page)    
            
        print(f'Done: {output_file}')
        cv.close()
    except Exception as e:
        print(f'An error occurred: {str(e)}')
        cv.close()


def pdf_to_xlsx(pdf_file, xlsx_file, page_range):
    xlsx_file += ".xlsx"
    if page_range is None:   
        pass
    elif len(page_range) == 2:
        page_range = [int(page) for page in page_range]
        page_range = list(range(page_range[0], page_range[1] + 1)) 
    else:          
        page_range = [int(page) + 1 for page in page_range]
               
    with pdfplumber.open(pdf_file) as pdf:
        all_data = []
        for page_number, page in enumerate(pdf.pages):
            if page_range is not None and page_number + 1 not in page_range:
                continue
            tables = page.extract_tables()
            for table_number, table in enumerate(tables):
                if not table:
                    continue
                
                df = pd.DataFrame(table)
                
                if df.empty:
                    continue

                all_data.append(df)

        if page_range is None and all_data:
            with pd.ExcelWriter(xlsx_file, engine='xlsxwriter') as writer:
                for i, df in enumerate(all_data):
                    df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False, header=False)
            print(f'Done: {xlsx_file}')
            
        elif page_range is not None and all_data:
            with pd.ExcelWriter(xlsx_file, engine='xlsxwriter') as writer:
                for i, df in enumerate(all_data):
                    df.to_excel(writer, sheet_name=f"Page_{page_range[i]}_Table_{i+1}", index=False, header=False)
            print(f'Done: {xlsx_file}')
        else:
            print("No data found.")
  

def pdf_to_pptx(pdf_file, pptx_file, page_range):
    pptx_file += ".pptx"
    if page_range is None:   
        pass
    elif len(page_range) == 2:
        page_range = [int(page) for page in page_range]
        page_range = list(range(page_range[0], page_range[1] + 1)) 
    else:          
        page_range = [int(page) + 1 for page in page_range]
        
    prs = Presentation()

    with fitz.open(pdf_file) as pdf_document:
        for page_number, pdf_page in enumerate(pdf_document):
            if page_range is not None and page_number + 1 not in page_range:
                continue

            pdf_page_rect = pdf_page.rect
            pdf_page_width = pdf_page_rect[2] - pdf_page_rect[0]
            pdf_page_height = pdf_page_rect[3] - pdf_page_rect[1]

            slide = prs.slides.add_slide(prs.slide_layouts[6]) 
            slide_width = Inches(pdf_page_width / 72.0)
            slide_height = Inches(pdf_page_height / 72.0)
            prs.slide_width = slide_width
            prs.slide_height = slide_height

            img = pdf_page.get_pixmap(matrix=fitz.Matrix(300/72, 300/72))
            img_path = f"temp_image.png"
            img.save(img_path)

            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(img_path, left, top, width=slide_width, height=slide_height)

            os.remove(img_path)

    prs.save(pptx_file)
    print(f'Done: {pptx_file}')


def pdf_to_images(pdf_file, image_folder, page_range=None):
    if not os.path.exists(image_folder):
        os.mkdir(image_folder)

    with fitz.open(pdf_file) as pdf_document:
        num_pages = len(pdf_document)

        if page_range is None:
            page_range = range(1, num_pages)
        elif len(page_range) == 2:
            page_range = range(int(page_range[0]), int(page_range[1]))
          

        for page_number in page_range:
            pdf_page = pdf_document[int(page_number)]

            img = pdf_page.get_pixmap(matrix=fitz.Matrix(300/72, 300/72))
            img_bytes = img.samples

            image = Image.frombytes("RGB", [img.width, img.height], img_bytes)
            image_file = os.path.join(image_folder, f"page_{int(page_number)+1}.png")
            image.save(image_file, "PNG")

    print(f'Done: {image_folder}')
        

def convert_file(pdf_file, output_file, choice_type, page, doc_type):
    if choice_type == "1":
        pdf_to_docx(pdf_file, output_file + "." + doc_type, page)
    elif choice_type == "2":
        pdf_to_xlsx(pdf_file, output_file, page)
    elif choice_type == "3":
        pdf_to_pptx(pdf_file, output_file, page)
    elif choice_type == "4":
        pdf_to_images(pdf_file, output_file, page)
             

def choice_page(number_of_pages):
    start_page = None
    end_page = None
    page = None
    if number_of_pages == "1":
        start_page = str(int(input("Start: ")) - 1)
        end_page = str(input("End: "))
        page = [start_page, end_page]
        return page
    elif number_of_pages == "2":
        page = [str(int(input("Page: ")) - 1)] 
        return page
    else:
        return page
    

def output_file_preparation(pdf_file, output_file):
    parts = pdf_file.split('/') 
    file_name = parts[-1].split(".")
    if output_file == "":
        output_file += file_name[0]
    else:
        output_file += "/" + file_name[0]
        
    return output_file


def main():
    pdf_file = input("Enter path and file name: ")
    choice_type = input("Choose an option:\n1. Convert to DOCX/DOC\n2. Convert to XLSX\n3. Convert to PPTX\n4. Convert to Images\n")
    doc_type = None
    if choice_type == "1":
        doc_type = input("Enter doc/docx: ")
    number_of_pages = input("Choose an option:\n1. Page range\n2. Only one page\n3. All file\n")
    
    page = choice_page(number_of_pages)
        
    output_file = input("Enter the path where you want to save the file: ") 
    output_file = output_file_preparation(pdf_file, output_file)
    
    convert_file(pdf_file, output_file, choice_type, page, doc_type)
    

if __name__ == "__main__":
    main()