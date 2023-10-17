from pdf2docx import Converter
import pdfplumber
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import os
import fitz
from PIL import Image


def pdf_to_docx(pdf_file, output_file, page):
    cv = Converter(pdf_file)
    try:
        if page is None:   
            start_page, end_page = 0, None
            cv.convert(output_file, start=start_page, end=end_page)
        elif len(page) == 2:
            start_page, end_page = page[0], page[1]
            cv.convert(output_file, start=start_page, end=end_page)
        else:          
            cv.convert(output_file, pages=page)    
            
        print(f'Done: {output_file}')
        cv.close()
    except Exception as e:
        print(f'An error occurred: {str(e)}')
        cv.close()


def pdf_to_xlsx(pdf_file, xlsx_file, page_range):
    if page_range is None:   
        pass
    elif len(page_range) == 2:
        page_range = [int(page) for page in page_range]
        page_range = list(range(page_range[0], page_range[1])) 
    else:          
        page_range = [int(page) for page in page_range]
        print(page_range)
    with pdfplumber.open(pdf_file) as pdf:
        all_data = []
        for page_number, page in enumerate(pdf.pages):
            if page_range is not None and page_number not in page_range:
                continue
            tables = page.extract_tables()
            for table_number, table in enumerate(tables):
                if not table:
                    continue
                
                df = pd.DataFrame(table)
                
                if df.empty:
                    continue

                all_data.append(df)

        if page_range is None and all_data:
            with pd.ExcelWriter(xlsx_file, engine='xlsxwriter') as writer:
                for i, df in enumerate(all_data):
                    df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False, header=False)
            print(f'Done: {xlsx_file}')
            
        elif page_range is not None and all_data:
            with pd.ExcelWriter(xlsx_file, engine='xlsxwriter') as writer:
                for i, df in enumerate(all_data):
                    df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False, header=False)
            print(f'Done: {xlsx_file}')
        else:
            print("No data found.")
  

def pdf_to_pptx(pdf_file, pptx_file, page_range):
    if page_range is None:   
        pass
    elif len(page_range) == 2:
        page_range = [int(page) + 1 for page in page_range]
        page_range = list(range(page_range[0], page_range[1])) 
    else:          
        page_range = [int(page) + 1 for page in page_range]
        
    prs = Presentation()

    with fitz.open(pdf_file) as pdf_document:
        for page_number, pdf_page in enumerate(pdf_document):
            if page_range is not None and page_number + 1 not in page_range:
                continue

            pdf_page_rect = pdf_page.rect
            pdf_page_width = pdf_page_rect[2] - pdf_page_rect[0]
            pdf_page_height = pdf_page_rect[3] - pdf_page_rect[1]

            slide = prs.slides.add_slide(prs.slide_layouts[6]) 
            slide_width = Inches(pdf_page_width / 72.0)
            slide_height = Inches(pdf_page_height / 72.0)
            prs.slide_width = slide_width
            prs.slide_height = slide_height

            img = pdf_page.get_pixmap(matrix=fitz.Matrix(300/72, 300/72))
            img_path = f"temp_image.png"
            img.save(img_path)

            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(img_path, left, top, width=slide_width, height=slide_height)

            os.remove(img_path)

    prs.save(pptx_file)
    print(f'Done: {pptx_file}')


def pdf_to_images(pdf_file, image_folder, page_range=None):
    if not os.path.exists(image_folder):
        os.mkdir(image_folder)

    with fitz.open(pdf_file) as pdf_document:
        num_pages = len(pdf_document)

        if page_range is None:
            page_range = range(0, num_pages)
        elif len(page_range) == 2:
            page_range = range(int(page_range[0]), int(page_range[1]))

        for page_number in page_range:
            pdf_page = pdf_document[int(page_number)]

            img = pdf_page.get_pixmap(matrix=fitz.Matrix(300/72, 300/72))
            img_bytes = img.samples

            image = Image.frombytes("RGB", [img.width, img.height], img_bytes)
            image_file = os.path.join(image_folder, f"page_{int(page_number)+1}.jpg")
            image.save(image_file, "JPEG")

    print(f'Done: {image_folder}')
        

def convert_file(pdf_file, output_file, page):
    file_extension = output_file.split(".")[-1]
    
    if file_extension == "doc" or file_extension == "docx":
        pdf_to_docx(pdf_file, output_file, page)
    elif file_extension == "xlsx":
        pdf_to_xlsx(pdf_file, output_file, page)
    elif file_extension == "ppt":
        pdf_to_pptx(pdf_file, output_file, page)
    elif file_extension == "jpg":
        pdf_to_images(pdf_file, output_file, page)
             

def main():  
    text_input = input().split()  
    pdf_file = text_input[0]
    output_file = text_input[1]  
    if text_input[2] == "all":
        page = None
    elif len(text_input[2].split("-")) == 1:
         page = [str(int(text_input[2])-1)]
    elif len(text_input[2].split("-")) == 2:
        page = text_input[2].split("-")
        page[0] = str(int(page[0])-1)
    convert_file(pdf_file, output_file, page)
    

if __name__ == "__main__":
    main()
    